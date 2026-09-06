# Copyright 2026 National Payments Corporation of India
# SPDX-License-Identifier: MIT

"""Shared helpers for the per-domain dashboard routers.

Holds the bits used by more than one domain module: the cert lifecycle order
(progress + certification) and the markdown→native-format converters (the
download endpoints in `changes`). Domain-specific helpers live in their own
module.
"""
import io

# Cert lifecycle order — shared by `progress` (auto-advance) and `certification`.
_CERT_STATUS_ORDER = ["received", "deployed", "tested", "ready_for_certification"]


# ── Markdown → native-format converters ───────────────────────────────────────
# Used by the download endpoints. NPCI generates a .docx for every Product Kit
# doc type and an additional .pptx for product_deck. Some inbound A2A
# change_communication payloads don't carry the binary bytes (path-resolution
# bug on the sender, or files not yet generated). In those cases the partner
# platform always has the markdown content — we synthesise the binary on-the-fly
# so every doc remains downloadable in the format NPCI advertises.

def markdown_to_docx_bytes(content: str, title: str = "") -> bytes:
    """Render markdown content into a .docx. Honours h1-h3, bullet
    lists, numbered lists, and paragraphs. Tables and inline
    formatting fall through as plain text — good enough for
    rollout documents which are mostly headings + prose."""
    from docx import Document
    doc = Document()
    if title:
        doc.add_heading(title, level=0)
    for raw in (content or "").splitlines():
        line = raw.rstrip()
        if not line:
            doc.add_paragraph("")
            continue
        if line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith(("- ", "* ")):
            doc.add_paragraph(line[2:], style="List Bullet")
        elif len(line) > 2 and line[0].isdigit() and line[1:3] in (". ", ".\t"):
            doc.add_paragraph(line.split(".", 1)[1].strip(), style="List Number")
        else:
            doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def markdown_to_pptx_bytes(content: str, title: str = "") -> bytes:
    """Render markdown into a basic .pptx — used only for product_deck.
    Every `# Heading` becomes a new slide; bullets and paragraphs
    under it land in the body placeholder. Title slide carries the
    doc title.
    """
    from pptx import Presentation
    prs = Presentation()
    # Title slide
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title or "Product Deck"
    if len(s.placeholders) > 1:
        s.placeholders[1].text = "Generated from rollout content"

    cur_title = None
    cur_body: list[str] = []

    def flush():
        if cur_title is None and not cur_body:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = cur_title or "Section"
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = cur_body[0] if cur_body else ""
            for line in cur_body[1:]:
                p = tf.add_paragraph()
                p.text = line

    for raw in (content or "").splitlines():
        line = raw.rstrip()
        if line.startswith("# "):
            flush()
            cur_title = line[2:]
            cur_body = []
        elif line.startswith(("## ", "### ")):
            cur_body.append(line.lstrip("#").strip())
        elif line.startswith(("- ", "* ")):
            cur_body.append("• " + line[2:])
        elif line.strip():
            cur_body.append(line)
    flush()

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _doc_title(doc) -> str:
    """Human-readable title from a doc_type string."""
    return (doc.doc_type or "Document").replace("_", " ").title()
